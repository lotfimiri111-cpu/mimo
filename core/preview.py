"""
core/preview.py v5 — قراءة ألوان PPTX من XML مباشرةً
"""
import base64, io, logging, os, threading
from typing import List, Optional

log = logging.getLogger(__name__)
_cache: dict = {}
_cache_lock = threading.Lock()
MAX_SLIDES = 3

CAIRO_PATHS = [
    os.path.expanduser("~/.fonts/cairo/Cairo.ttf"),
    "/root/.fonts/cairo/Cairo.ttf",
    "/tmp/fonts/cairo/Cairo.ttf",
    "/home/user/.fonts/cairo/Cairo.ttf",
    "/opt/render/project/src/.fonts/cairo/Cairo.ttf",
    "/opt/render/.fonts/cairo/Cairo.ttf",
    # Amiri fallback
    os.path.expanduser("~/.fonts/cairo/Amiri-Regular.ttf"),
    "/root/.fonts/cairo/Amiri-Regular.ttf",
    "/tmp/fonts/cairo/Amiri-Regular.ttf",
]
FALLBACK_FONTS = [
    "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
    "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf",
    "/usr/share/fonts/truetype/freefont/FreeSansBold.ttf",
    "/usr/share/fonts/truetype/freefont/FreeSans.ttf",
]

def _find_font(size=16, bold=False):
    from PIL import ImageFont
    # Cairo أولاً
    for p in CAIRO_PATHS:
        if os.path.exists(p):
            try: return ImageFont.truetype(p, size)
            except: pass
    # Fallback
    for p in FALLBACK_FONTS:
        if os.path.exists(p):
            try: return ImageFont.truetype(p, size)
            except: pass
    return ImageFont.load_default()


def get_cached_preview(pid):
    with _cache_lock: return _cache.get(pid)

def set_cached_preview(pid, slides):
    with _cache_lock: _cache[pid] = slides


def pptx_to_preview_images(pptx_path: str, watermark: bool = True) -> List[str]:
    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
        W = max(960, int((prs.slide_width  or 9144000) / 9525))
        H = max(540, int((prs.slide_height or 5143500) / 9525))
        results = []
        for slide in list(prs.slides)[:MAX_SLIDES]:
            img = _render_slide(slide, W, H)
            if watermark:
                img = _add_watermark(img)
            buf = io.BytesIO()
            img.save(buf, "JPEG", quality=90)
            results.append(base64.b64encode(buf.getvalue()).decode())
        return results
    except Exception as e:
        log.error(f"preview v5 failed: {e}", exc_info=True)
        return []


def _hex(h) -> tuple:
    h = str(h).lstrip("#")
    if len(h) == 6:
        return (int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))
    return (200,200,200)


def _read_color_xml(elem) -> Optional[tuple]:
    """قراءة لون من عنصر XML بشكل مباشر"""
    if elem is None: return None
    try:
        from lxml import etree
        # srgbClr
        for child in elem.iter():
            tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
            if tag == 'srgbClr':
                val = child.get('val','')
                if val: return _hex(val)
            elif tag == 'sysClr':
                val = child.get('lastClr','')
                if val: return _hex(val)
            elif tag == 'prstClr':
                # preset colors
                prstMap = {'white':(255,255,255),'black':(0,0,0),'red':(255,0,0),
                           'blue':(0,0,255),'yellow':(255,255,0),'green':(0,128,0)}
                return prstMap.get(child.get('val',''), (128,128,128))
    except: pass
    return None


def _bg_color(slide) -> tuple:
    """استخراج لون خلفية الشريحة"""
    try:
        # محاولة 1: pptx API
        fill = slide.background.fill
        if fill.type is not None:
            c = fill.fore_color.rgb
            return (c.r, c.g, c.b)
    except: pass
    try:
        # محاولة 2: XML مباشر
        bg_elem = slide.background._element
        c = _read_color_xml(bg_elem)
        if c: return c
    except: pass
    return (255, 255, 255)


def _shape_fill_color(shape) -> Optional[tuple]:
    try:
        fill = shape.fill
        from pptx.enum.dml import MSO_FILL
        # BACKGROUND(5) = شفاف — لا نرسم خلفية
        if hasattr(MSO_FILL, 'BACKGROUND') and fill.type == MSO_FILL.BACKGROUND:
            return None
        if fill.type is not None and fill.type != 5:
            try:
                c = fill.fore_color.rgb
                return (c.r, c.g, c.b)
            except: pass
    except: pass
    try:
        ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        sp_elem = shape._element
        solid = sp_elem.find(f'.//{{{ns}}}solidFill')
        if solid is not None:
            c = _read_color_xml(solid)
            if c: return c
    except: pass
    return None


def _text_color(run) -> tuple:
    try:
        rgb = run.font.color.rgb
        return (rgb.r, rgb.g, rgb.b)
    except: pass
    try:
        r_elem = run._r
        c = _read_color_xml(r_elem)
        if c: return c
    except: pass
    return None

def _text_color_from_xml(run) -> Optional[tuple]:
    """قراءة لون النص من XML مباشرةً — أكثر موثوقية"""
    # محاولة 1: pptx API
    try:
        rgb = run.font.color.rgb
        return (rgb.r, rgb.g, rgb.b)
    except: pass
    # محاولة 2: XML مباشر
    try:
        # ابحث في عنصر الـ run عن lون
        ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        r_elem = run._r
        # البحث في rPr عن solidFill/srgbClr
        rPr = r_elem.find(f'{{{ns}}}rPr')
        if rPr is not None:
            solid = rPr.find(f'{{{ns}}}solidFill')
            if solid is not None:
                c = _read_color_xml(solid)
                if c: return c
    except: pass
    # محاولة 3: parent paragraph
    try:
        para_elem = run._r.getparent()
        ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        pPr = para_elem.find(f'{{{ns}}}pPr')
        if pPr is not None:
            c = _read_color_xml(pPr)
            if c: return c
    except: pass
    return None


def _smart_text_color(bg: tuple) -> tuple:
    """اختر لون النص تلقائياً بحسب لمعان الخلفية"""
    lum = 0.299*bg[0] + 0.587*bg[1] + 0.114*bg[2]
    return (255,255,255) if lum < 128 else (20,20,20)


def _render_slide(slide, W, H):
    from PIL import Image, ImageDraw
    bg = _bg_color(slide)
    img = _draw_bg(bg, W, H, slide)
    draw = ImageDraw.Draw(img)
    for shape in slide.shapes:
        try: _draw_shape(draw, img, shape, W, H, bg)
        except: pass
    return img


def _draw_bg(bg_color, W, H, slide):
    """رسم الخلفية مع gradient إذا وُجد"""
    from PIL import Image
    try:
        # قراءة الـ gradient من XML
        from lxml import etree
        ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        bg_elem = slide.background._element
        grad = bg_elem.find(f'.//{{{ns}}}gradFill')
        if grad is not None:
            stops = []
            for gs in grad.findall(f'.//{{{ns}}}gs'):
                pos = int(gs.get('pos','0')) / 100000.0
                c = _read_color_xml(gs)
                if c: stops.append((pos, c))
            if len(stops) >= 2:
                return _gradient_img(W, H, sorted(stops, key=lambda x: x[0]))
    except: pass
    # لون صلب مع تدرج خفيف
    return _soft_gradient(bg_color, W, H)


def _soft_gradient(color, W, H):
    """تدرج خفيف من اللون نفسه"""
    from PIL import Image
    img = Image.new("RGB", (W, H))
    px = img.load()
    r,g,b = color
    for y in range(H):
        t = y/H
        # أفتح قليلاً في الأعلى
        factor = 1.0 + 0.25*(1-t) - 0.1*t
        nr = min(255, int(r*factor))
        ng = min(255, int(g*factor))
        nb = min(255, int(b*factor))
        for x in range(W): px[x,y] = (nr,ng,nb)
    return img


def _gradient_img(W, H, stops):
    from PIL import Image
    img = Image.new("RGB", (W, H))
    px = img.load()
    for y in range(H):
        t = y/H
        c1,c2 = stops[0][1], stops[-1][1]
        for i in range(len(stops)-1):
            p0,col0 = stops[i]; p1,col1 = stops[i+1]
            if p0 <= t <= p1:
                lt = (t-p0)/(p1-p0) if p1>p0 else 0
                c1 = tuple(int(col0[j]+(col1[j]-col0[j])*lt) for j in range(3))
                break
        for x in range(W): px[x,y] = c1
    return img


def _draw_shape(draw, img, shape, W, H, slide_bg):
    from PIL import ImageFont
    from pptx.enum.dml import MSO_THEME_COLOR
    EMU = 914400.0; DPI = 96
    L = int((shape.left  or 0)/EMU*DPI)
    T = int((shape.top   or 0)/EMU*DPI)
    SW = int((shape.width or 0)/EMU*DPI)
    SH = int((shape.height or 0)/EMU*DPI)

    # خلفية الشكل — نتجاهل BACKGROUND(5) لأنه شفاف
    fc = _shape_fill_color(shape)
    if fc and SW > 0 and SH > 0:
        draw.rectangle([L,T,L+SW,T+SH], fill=fc)

    if not shape.has_text_frame: return

    # لون افتراضي للنص: نستخدم خلفية الشريحة لتحديد التباين
    # لكن إذا كان الشكل له خلفية — نحسب بالنسبة لها
    bg_for_contrast = fc if fc else slide_bg
    default_color = _smart_text_color(bg_for_contrast)

    y = T + 6
    for para in shape.text_frame.paragraphs:
        text = para.text.strip()
        if not text: y += 8; continue
        # حجم ولون من الـ runs — نقرأ اللون المُحدَّد مباشرةً
        fs, tc, bold = 14, None, False
        for run in para.runs:
            try:
                if run.font.size: fs = max(8, min(int(run.font.size/12700), 60))
                if run.font.bold: bold = True
                # نقرأ اللون من الـ XML مباشرة لأن pptx API أحياناً يفشل
                c = _text_color_from_xml(run)
                if c: tc = c
            except: pass
            break
        if tc is None: tc = default_color
        font = _find_font(fs, bold)
        lines = _wrap(text, font, SW-12)
        for line in lines:
            if y >= T+SH+10: break
            draw.text((L+6, y), line, fill=tc, font=font)
            y += fs+4
        if y >= T+SH+10: break


def _wrap(text, font, max_w):
    from PIL import ImageDraw, Image
    if max_w <= 20: return [text]
    d = ImageDraw.Draw(Image.new("RGB",(1,1)))
    def tw(t):
        try: bb=d.textbbox((0,0),t,font=font); return bb[2]-bb[0]
        except: return len(t)*8
    if tw(text) <= max_w: return [text]
    words = text.split()
    lines, cur = [], ""
    for w in words:
        test = (cur+" "+w).strip()
        if tw(test) <= max_w: cur=test
        else:
            if cur: lines.append(cur)
            cur=w
    if cur: lines.append(cur)
    return lines or [text]


def _add_watermark(img):
    from PIL import Image, ImageDraw
    W, H = img.size
    overlay = Image.new("RGBA",(W,H),(0,0,0,0))
    d = ImageDraw.Draw(overlay)
    font = _find_font(max(22, W//26), bold=True)
    text = "مذكرتي Pro — معاينة فقط"
    try:
        bb = d.textbbox((0,0),text,font=font)
        tw,th = bb[2]-bb[0], bb[3]-bb[1]
    except:
        tw,th = len(text)*14, 28
    pad=40
    tl = Image.new("RGBA",(tw+pad,th+pad),(0,0,0,0))
    td = ImageDraw.Draw(tl)
    td.text((pad//2,pad//2), text, fill=(255,255,255,80), font=font)
    rot = tl.rotate(-30, expand=True)
    rw,rh = rot.size
    sx = max(rw+30, W//3)
    sy = max(rh+20, H//3)
    for row in range(-1, H//sy+2):
        for col in range(-1, W//sx+2):
            overlay.paste(rot, (col*sx-rw//2, row*sy-rh//2), rot)
    out = Image.alpha_composite(img.convert("RGBA"), overlay).convert("RGB")
    return out
